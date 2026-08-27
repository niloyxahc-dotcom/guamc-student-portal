from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 Widescreen
prs.slide_height = Inches(7.5)

# Color Palette (Medical & Modern Theme)
COLOR_PRIMARY = RGBColor(16, 44, 87)     # Deep Navy
COLOR_SECONDARY = RGBColor(53, 162, 159) # Teal / Cyan Accent
COLOR_BG_DARK = RGBColor(245, 247, 250)  # Light Soft Grey/Blue
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)  # Charcoal
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_MUTED = RGBColor(100, 116, 139)

def apply_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_DARK

def create_header(slide, title_text, category_text="ইউনানী শারীরতত্ত্ব ও ইলমুল আদভিয়া"):
    # Header Banner shape
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_PRIMARY
    banner.line.color.rgb = COLOR_PRIMARY

    # Accent Stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_SECONDARY
    stripe.line.color.rgb = COLOR_SECONDARY

    # Title Text
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.15)

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

def add_content_box(slide, left, top, width, height, title, points, bg_white=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE if bg_white else COLOR_PRIMARY
    shape.line.color.rgb = COLOR_SECONDARY if bg_white else COLOR_PRIMARY

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.25)

    p_header = tf.paragraphs[0]
    p_header.text = title
    p_header.font.size = Pt(15)
    p_header.font.bold = True
    p_header.font.color.rgb = COLOR_PRIMARY if bg_white else COLOR_WHITE

    for pt in points:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN if bg_white else COLOR_WHITE
        p.space_before = Pt(6)

# ==================== SLIDE 1: Title Slide ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
card.fill.solid()
card.fill.fore_color.rgb = COLOR_WHITE
card.line.color.rgb = COLOR_SECONDARY

tf1 = card.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.8)
tf1.margin_top = Inches(0.6)

p1 = tf1.paragraphs[0]
p1.text = "ইউনানী শারীরতত্ত্ব ও ইলমুল আদভিয়ার মেলবন্ধন"
p1.font.size = Pt(26)
p1.font.bold = True
p1.font.color.rgb = COLOR_PRIMARY

p2 = tf1.add_paragraph()
p2.text = "আরওয়াহ (Arwah), কোয়াহ (Quwa) ও আফ’আল (Af’al)-এর গভীর বিশ্লেষণ এবং আধুনিক চিকিৎসাবিজ্ঞানের তুলনামূলক প্রেক্ষাপট"
p2.font.size = Pt(14)
p2.font.color.rgb = COLOR_MUTED
p2.space_before = Pt(10)

p3 = tf1.add_paragraph()
p3.text = "স্পিকার: Dr. A M Zahid Hasan, BCS (Health)"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = COLOR_SECONDARY
p3.space_before = Pt(35)

p4 = tf1.add_paragraph()
p4.text = "ইউনানী চিকিৎসাবিজ্ঞান বিভাগ | বিষয়: ইলমুল আদভিয়া"
p4.font.size = Pt(13)
p4.font.color.rgb = COLOR_TEXT_MAIN
p4.space_before = Pt(5)

# ==================== SLIDE 2: উমূরে তবাইয়াহ ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide2)
create_header(slide2, "উমূরে তবাইয়াহ (Umur-e-Tabi'yah) — ৭টি মৌলিক ভিত্তি")
add_content_box(slide2, 0.8, 1.5, 5.6, 5.4, "মৌলিক চার ভিত্তি (Structural Basis)", [
    "আর্কান (Arkan): প্রাথমিক মৌলিক উপাদান (আগুন, পানি, বাতাস, মাটি)।",
    "মিজাজ (Mizaj): মেজাজ বা প্রাকৃতিক ভারসাম্য (Temperament)।",
    "আখলাত (Akhlat): চার শারীরিক রস (দম/রক্ত, বলগম/কফ, সাফরা/পিত্ত, সৌদা)।",
    "আ'জা (A'za): শারীরিক অঙ্গপ্রত্যঙ্গ ও কোষকলা (Organs & Tissues)।"
])
add_content_box(slide2, 6.9, 1.5, 5.6, 5.4, "শারীরবৃত্তীয় তিন চালিকাশক্তি (Functional Triad)", [
    "আরওয়াহ (Arwah): জীবনীশক্তি বহনকারী সূক্ষ্ম বাষ্পীয় মাধ্যম (Pneuma)।",
    "কোয়াহ (Quwa): অঙ্গপ্রত্যঙ্গের অন্তর্নিহিত কার্যক্ষমতা বা অনুষদ (Faculties)।",
    "আফ'আল (Af'al): ক্ষমতার চূড়ান্ত বহিঃপ্রকাশ বা শারীরিক কাজ (Functions)।",
    "ক্লিনিক্যাল গুরুত্ব: রোগমুক্তির জন্য এই সাতটি উপাদানের সমতা রক্ষা অপরিহার্য।"
])

# ==================== SLIDE 3: ইলমুল আদভিয়া ও ফার্মাকোডাইনামিক্স ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide3)
create_header(slide3, "ইলমুল আদভিয়া ও ফার্মাকোডাইনামিক্সের সমান্তরাল দর্শন")
add_content_box(slide3, 0.8, 1.5, 5.6, 5.4, "ইউনানী ড্রাগ অ্যাকশন মেকানিজম", [
    "ভেষজ ওষুধ রোগীর দেহে মিজাজের ভারসাম্যহীনতা সংশোধন করে।",
    "ওষুধ প্রথমে আখলাত (রসের মান ও পরিমাণ)-এ রূপান্তর ঘটায়।",
    "আখলাতের পরিশুদ্ধতা আরওয়াহ ও কোয়াহকে সক্রিয় করে তোলে।",
    "চূড়ান্ত পর্যায়ে স্বাভাবিক ও সুস্থ আফ’আল পুনরুদ্ধার হয়।"
])
add_content_box(slide3, 6.9, 1.5, 5.6, 5.4, "আধুনিক ফার্মাকোলজি ও হোমিওস্টেসিস", [
    "ড্রাগ মলিকিউল নির্দিষ্ট রিসেপ্টরের সাথে বাইন্ড করে সংকেত দেয়।",
    "সেলুলার সিগন্যালিংয়ের মাধ্যমে এনজাইম ও ট্রান্সপোর্টার সক্রিয় হয়।",
    "কোষীয় মেটাবলিজম বাড়িয়ে হোমিওস্টেসিস বজায় রাখা হয়।",
    "ইউনানীর 'তবাইয়াত' এবং আধুনিক 'Milieu Intérieur' একই আত্মনিয়ন্ত্রণ রূপ।"
])

# ==================== SLIDE 4: আরওয়াহ — ধারণা ও উপাদান ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide4)
create_header(slide4, "আরওয়াহ (Arwah / Spirits) — ধারণা ও উপাদান")
add_content_box(slide4, 0.8, 1.5, 5.6, 5.4, "আরওয়াহর ক্লাসিক্যাল ধারণা", [
    "আরওয়াহ হলো 'রূহ'-এর বহুবচন; এটি কোনো অদৃশ্য বা অলৌকিক সত্তা নয়।",
    "এটি বাতাস ও রক্তের সূক্ষ্ম বাষ্পের সমন্বয়ে গঠিত জীবনী পরিবাহক।",
    "উপাদান ১: বাহ্যিক বাতাস বা শ্বাসবায়ু (Inhaled Air / O2)।",
    "উপাদান ২: আখলাত-ই-লাতিফাহ (রক্তের পরিশোধিত নির্যাস)।"
])
add_content_box(slide4, 6.9, 1.5, 5.6, 5.4, "আধুনিক বায়োকেমিক্যাল তাৎপর্য", [
    "ধমনিস্থ অক্সিজেন সংবহন (PaO2 ও Hemoglobin Saturation)।",
    "কোষীয় শ্বসন (Cellular Respiration) ও গ্লুকোজ জারণ।",
    "মাইটোকন্ড্রিয়ায় উৎপন্ন এটিপি (ATP - Adenosine Triphosphate)।",
    "আরওয়াহ সরাসরি কোষের ভাইটালিটি এবং অক্সিজেনেশনের সমতুল্য।"
])

# ==================== SLIDE 5: আরওয়াহর প্রকারভেদ ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide5)
create_header(slide5, "আরওয়াহর প্রকারভেদ ও শারীরবৃত্তীয় কেন্দ্র")
add_content_box(slide5, 0.8, 1.5, 3.6, 5.4, "রূহ-ই-তবাইয়াহ (Natural)", [
    "কেন্দ্র: যকৃৎ (Liver)।",
    "বাহক: শিরাতন্ত্র (Venous System)।",
    "কাজ: পরিপাককৃত খাদ্যের পুষ্টি উপাদান সারা শরীরের টিস্যুতে পৌঁছে দেওয়া।"
])
add_content_box(slide5, 4.8, 1.5, 3.6, 5.4, "রূহ-ই-হায়ওয়ানিয়াহ (Vital)", [
    "কেন্দ্র: হৃৎপিণ্ড (Heart)।",
    "বাহক: ধমনি ব্যবস্থা (Arterial System)।",
    "কাজ: সমগ্র দেহে প্রাণশক্তি, জীবনীপ্রবাহ এবং স্বাভাবিক শারীরিক তাপমাত্রা বজায় রাখা।"
])
add_content_box(slide5, 8.8, 1.5, 3.6, 5.4, "রূহ-ই-নাফসানিয়াহ (Psychic)", [
    "কেন্দ্র: মস্তিষ্ক (Brain)।",
    "বাহক: স্নায়ুতন্ত্র (Nerves)।",
    "কাজ: চেতনা, ইন্দ্রিয় সংবেদন (Sensation) এবং পেশির ঐচ্ছিক সঞ্চালন পরিচালনা করা।"
])

# ==================== SLIDE 6: আরওয়াহ বনাম আধুনিক সার্কুলেশন ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide6)
create_header(slide6, "আরওয়াহ বনাম আধুনিক সিস্টেমিক সার্কুলেশন")
add_content_box(slide6, 0.8, 1.5, 11.7, 5.4, "শারীরবৃত্তীয় সিস্টেমের সরাসরি তুলনা", [
    "রূহ-ই-তবাইয়াহ $\leftrightarrow$ হেপাটিক পোর্টাল সার্কুলেশন, লিভার মেটাবলিজম ও প্লাজমা প্রোটিন পরিবহন।",
    "রূহ-ই-হায়ওয়ানিয়াহ $\leftrightarrow$ পালমোনারি গ্যাস এক্সচেঞ্জ, মায়োকার্ডিয়াল পারফিউশন ও সিস্টেমিক ব্লাড প্রেসার।",
    "রূহ-ই-নাফসানিয়াহ $\leftrightarrow$ সেরিব্রাল পারফিউশন, ব্লাড-ব্রেন ব্যারিয়ার (BBB) পারাপার ও অ্যাকশন পটেনশিয়াল।"
])

# ==================== SLIDE 7: কোয়াহ — শক্তি ও অনুষদ ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide7)
create_header(slide7, "কোয়াহ (Quwa / Faculties) — শক্তি ও অনুষদ")
add_content_box(slide7, 0.8, 1.5, 5.6, 5.4, "কোয়াহর সংজ্ঞা ও প্রকৃতি", [
    "শরীরের অঙ্গপ্রত্যঙ্গকে কার্যকর রাখার অন্তর্নিহিত প্রাকৃতিক সামর্থ্য।",
    "আরওয়াহ হলো এনার্জি বা জ্বালানি (Energy Carrier)।",
    "কোয়াহ হলো সেই এনার্জি ব্যবহারকারী সেলুলার মেকানিজম (Functional Engine)।"
])
add_content_box(slide7, 6.9, 1.5, 5.6, 5.4, "প্রধান তিনটি বিভাগ", [
    "কুওয়াত-ই-তবাইয়াহ (Natural Faculty): পুষ্টি, বৃদ্ধি ও বংশরক্ষা।",
    "কুওয়াত-ই-হায়ওয়ানিয়াহ (Vital Faculty): কার্ডিও-পালমোনারি প্রাণশক্তি।",
    "কুওয়াত-ই-নাফসানিয়াহ (Psychic Faculty): স্নায়বিক অনুভূতি ও মোটর শক্তি।"
])

# ==================== SLIDE 8: কুওয়াত-ই-তবাইয়াহ ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide8)
create_header(slide8, "কুওয়াত-ই-তবাইয়াহ (Natural Faculty) — পুষ্টি ও বৃদ্ধি")
add_content_box(slide8, 0.8, 1.5, 5.6, 5.4, "কুওয়াত-ই-গাযিয়াহ (Nutritive Power)", [
    "জাযিবাহ (Attractive): পুষ্টি শোষণ বা গ্রহণ (Cellular Uptake)।",
    "মাসিকাহ (Retentive): পরিপাক না হওয়া পর্যন্ত উপাদান ধরে রাখা।",
    "হাযিমাহ (Digestive): রাসায়নিক ভাঙন ও আত্মীকরণ (Metabolism)।",
    "দাফিআহ (Expulsive): মেটাবলিক বর্জ্য নিষ্কাশন (Excretion)।"
])
add_content_box(slide8, 6.9, 1.5, 5.6, 5.4, "বৃদ্ধি ও প্রজনন শক্তি", [
    "কুওয়াত-ই-নামিয়াহ (Growth Power): কোষ বিভাজন ও দেহের বৃদ্ধি (Cell Proliferation & GH)।",
    "কুওয়াত-ই-মুওয়াল্লিদাহ (Generative Power): স্পার্মাটোজেনেসিস, ওভুলেশন ও প্রজনন।"
])

# ==================== SLIDE 9: কুওয়াত-ই-হায়ওয়ানিয়াহ ====================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide9)
create_header(slide9, "কুওয়াত-ই-হায়ওয়ানিয়াহ (Vital Faculty) — জীবনী অনুষদ")
add_content_box(slide9, 0.8, 1.5, 5.6, 5.4, "ইউনানী দৃষ্টিভঙ্গি", [
    "হৃৎপিণ্ড ও রক্তনালিকে সচল রাখার মূল চালিকাশক্তি।",
    "হৃৎপিণ্ডের নিয়মিত সংকোচন (Systole) ও প্রসারণ (Diastole) নিয়ন্ত্রণ।",
    "ধমনিজুড়ে পালস ও জীবনীপ্রবাহ নিশ্চিত করা।"
])
add_content_box(slide9, 6.9, 1.5, 5.6, 5.4, "আধুনিক ফিজিওলজি কোরিলেশন", [
    "সাইনোঅ্যাট্রিয়াল নোড (SA Node)-এর পেসমেকার অটোমেটিসিটি।",
    "মায়োকার্ডিয়াল কন্ট্রাকটিলিটি ও কার্ডিয়াক আউটপুট।",
    "অটোনমিক নার্ভাস সিস্টেম (Sympathetic ও Parasympathetic টোন)।"
])

# ==================== SLIDE 10: কুওয়াত-ই-নাফসানিয়াহ ====================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide10)
create_header(slide10, "কুওয়াত-ই-নাফসানিয়াহ (Psychic Faculty) — স্নায়বিক ও মানসিক ক্ষমতা")
add_content_box(slide10, 0.8, 1.5, 5.6, 5.4, "কুওয়াত-ই-মুদরিকাহ (Perceptive)", [
    "বাহ্যিক ইন্দ্রিয়: পঞ্চেন্দ্রিয় (দর্শন, শ্রবণ, ঘ্রাণ, স্বাদ, স্পর্শ)।",
    "অভ্যন্তরীণ ইন্দ্রিয়: সাধারণ অনুভূতি (Hiss-e-Mushtarak), কল্পনা (Khayal), চিন্তা (Fikr) ও স্মৃতি (Hafiza)।"
])
add_content_box(slide10, 6.9, 1.5, 5.6, 5.4, "কুওয়াত-ই-মুহাররিকাহ (Motive)", [
    "স্নায়বিক সংকেতের মাধ্যমে মাংসপেশির সংকোচন ও অঙ্গ সঞ্চালন।",
    "মোটর নার্ভ পাথওয়ে ও নিউরো-মাসকুলার জাংশনের কার্যকারিতা।"
])

# ==================== SLIDE 11: আফ’আল ====================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide11)
create_header(slide11, "আফ’আল (Af’al / Functions) — চূড়ান্ত শারীরিক প্রকাশ")
add_content_box(slide11, 0.8, 1.5, 5.6, 5.4, "আফ’আল-ই-মুফরাদা (Simple Actions)", [
    "একক শক্তি বা উপাদানের মাধ্যমে সম্পন্ন কাজ।",
    "উদাহরণ: কেবল খাদ্য শোষণ, পিত্ত নিঃসরণ বা রক্তনালির সংকোচন।"
])
add_content_box(slide11, 6.9, 1.5, 5.6, 5.4, "আফ’আল-ই-মুরাক্কাবা (Compound Actions)", [
    "একাধিক শক্তি ও অঙ্গের সমন্বিত জটিল কাজ।",
    "উদাহরণ: ক্ষুধা অনুভব হওয়া, খাদ্য চর্বণ থেকে পরিপাক ও মলত্যাগ, শারীরিক চলন।"
])

# ==================== SLIDE 12: তুলনামূলক ছক ====================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide12)
create_header(slide12, "ইউনানী ও আধুনিক শারীরতত্ত্বের ট্রায়াড ম্যাপিং")

# Create Table
rows, cols = 4, 3
left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.733), Inches(4.8)
table_shape = slide12.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(3.2)
table.columns[1].width = Inches(4.2)
table.columns[2].width = Inches(4.333)

headers = ["পর্যায়", "ইউনানী ভূমিকা", "আধুনিক সমতুল্য শারীরতত্ত্ব"]
data = [
    ["আরওয়াহ (Arwah)", "সূক্ষ্ম বাষ্পীয় মাধ্যম ও জীবনী পরিবাহক", "ধমনিস্থ O2 সংবহন, হিমোগ্লোবিন ও মাইটোকন্ড্রিয়াল ATP"],
    ["কোয়াহ (Quwa)", "অঙ্গের অন্তর্নিহিত ক্ষমতা ও কর্মদক্ষতা", "রিসেপ্টর, এনজাইম অ্যাক্টিভিটি ও হরমোনাল রেগুলেশন"],
    ["আফ'আল (Af'al)", "চূড়ান্ত দৃশ্যমান শারীরিক ক্রিয়াকলাপ", "কার্ডিয়াক আউটপুট, পেরিস্টালসিস ও নিউরাল রিফ্লেক্স"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_WHITE

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE if row_idx % 2 == 0 else RGBColor(235, 240, 248)
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN

# ==================== SLIDE 13: ইলমুল আদভিয়ায় ড্রাগ ক্লাসিফিকেশন ====================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide13)
create_header(slide13, "ইলমুল আদভিয়ায় ড্রাগ ক্লাসিফিকেশন ও ট্রায়াড়ের প্রভাব")
add_content_box(slide13, 0.8, 1.5, 5.6, 2.5, "মুহাররিকাত (Stimulants)", [
    "আরওয়াহর গতি বাড়িয়ে কোয়াহকে উদ্দীপিত করে।",
    "উদাহরণ: এফিড্রিন/ক্যাফেইন সমৃদ্ধ ভেষজ।"
])
add_content_box(slide13, 6.9, 1.5, 5.6, 2.5, "মুসাক্কিনাত (Sedatives)", [
    "অতিরিক্ত সক্রিয় কোয়াহ-ই-নাফসানিয়াহকে প্রশমিত করে ব্যথা ও অস্থিরতা কমায়।"
])
add_content_box(slide13, 0.8, 4.3, 5.6, 2.5, "মুকাব্বিয়াত (Tonics)", [
    "কুওয়াত-ই-তবাইয়াহ ও লিভার মেটাবলিজম বাড়িয়ে সাধারণ শারীরিক শক্তি জোগায়।"
])
add_content_box(slide13, 6.9, 4.3, 5.6, 2.5, "মুফাররিহাত (Exhilarants)", [
    "রূহ-ই-হায়ওয়ানিয়াহ ও নাফসানিয়াহকে পুষ্টি দিয়ে মানসিক সতেজতা আনে।"
])

# ==================== SLIDE 14: ক্লিনিক্যাল কেস স্টাডি ====================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide14)
create_header(slide14, "ক্লিনিক্যাল কেস স্টাডি — ধারণাগত মেলবন্ধন")
add_content_box(slide14, 0.8, 1.5, 5.6, 5.4, "লক্ষণ: দীর্ঘমেয়াদি শারীরিক ক্লান্তি ও দুর্বলতা", [
    "ইউনানী ডায়াগনোসিস: আরওয়াহর হ্রাস এবং কুওয়াত-ই-তবাইয়াহর দুর্বলতা (Zo'f-e-A'za)।",
    "আধুনিক ডায়াগনোসিস: সেলুলার হাইপোক্সিয়া, রক্তস্বল্পতা ও এটিপি উৎপাদনের ঘাটতি।"
])
add_content_box(slide14, 6.9, 1.5, 5.6, 5.4, "চিকিৎসা কৌশল (Management)", [
    "ইলমুল আদভিয়া অনুযায়ী: মুকাব্বিয়াত ও মুফাররিহাত ড্রাগ প্রয়োগ।",
    "শারীরবৃত্তীয় ফল: আরওয়াহ পুনরুজ্জীবিত $\rightarrow$ কোয়াহ ও আফ’আল স্বাভাবিকীকরণ।"
])

# ==================== SLIDE 15: লেকচার সারসংক্ষেপ ====================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide15)
create_header(slide15, "লেকচার সারসংক্ষেপ (Key Takeaways)")
add_content_box(slide15, 0.8, 1.5, 11.7, 5.4, "মূল শিক্ষা ও উপলব্ধি", [
    "আরওয়াহ, কোয়াহ এবং আফ’আল হলো যথাক্রমে জ্বালানি $\rightarrow$ সামর্থ্য $\rightarrow$ কাজের অবিচ্ছেদ্য শারীরবৃত্তীয় ধারা।",
    "প্রাচীন ইউনানী চিকিৎসা দর্শন মূলত আধুনিক সেলুলার ফিজিওলজি ও মেটাবলিজমের একটি দূরদর্শী তাত্ত্বিক রূপ।",
    "ইলমুল আদভিয়ায় ওষুধ নির্বাচনের মূল লক্ষ্য হলো এই ট্রায়াড়ের ভারসাম্যহীনতা দূর করে রোগ নিরাময় করা।"
])

# ==================== SLIDE 16: প্রশ্নোত্তর ====================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
bg16 = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg16.fill.solid()
bg16.fill.fore_color.rgb = COLOR_PRIMARY
bg16.line.fill.background()

card16 = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(1.8), Inches(8.333), Inches(3.9))
card16.fill.solid()
card16.fill.fore_color.rgb = COLOR_WHITE
card16.line.color.rgb = COLOR_SECONDARY

tf16 = card16.text_frame
tf16.word_wrap = True
tf16.margin_top = Inches(0.8)

p_q1 = tf16.paragraphs[0]
p_q1.text = "প্রশ্নোত্তর ও উন্মুক্ত আলোচনা (Q&A)"
p_q1.alignment = PP_ALIGN.CENTER
p_q1.font.size = Pt(24)
p_q1.font.bold = True
p_q1.font.color.rgb = COLOR_PRIMARY

p_q2 = tf16.add_paragraph()
p_q2.text = "ইউনানী শারীরতত্ত্ব ও আধুনিক চিকিৎসাবিজ্ঞানের মেলবন্ধন বিষয়ে আপনাদের কোনো জিজ্ঞাসা থাকলে করতে পারেন।"
p_q2.alignment = PP_ALIGN.CENTER
p_q2.font.size = Pt(13)
p_q2.font.color.rgb = COLOR_MUTED
p_q2.space_before = Pt(15)

p_q3 = tf16.add_paragraph()
p_q3.text = "ধন্যবাদ!"
p_q3.alignment = PP_ALIGN.CENTER
p_q3.font.size = Pt(20)
p_q3.font.bold = True
p_q3.font.color.rgb = COLOR_SECONDARY
p_q3.space_before = Pt(20)

# Save Presentation
output_filename = "Ilmul_Advia_Lecture_Dr_Zahid.pptx"
prs.save(output_filename)
print(f"Presentation successfully saved as '{output_filename}'")